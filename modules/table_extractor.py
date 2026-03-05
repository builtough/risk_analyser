"""
Table Extractor — adapted from production extraction code.

Accepts uploaded file bytes (BytesIO) directly — no disk/config dependency.
Returns a list of table dicts, each carrying the DataFrame + metadata.
Also provides table_to_chunks() to convert tables into text chunks
so all existing search, analysis, and RAG features work on table data.

Supported: PDF, DOCX, PPTX, XLSX/XLS, CSV
"""

import io
import re
import pandas as pd
from typing import List, Dict, Any, Optional

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ── PDF extraction settings (two-pass: lattice + stream) ─────────────────────

_PDF_LATTICE = dict(
    vertical_strategy="lines",
    horizontal_strategy="lines",
    intersection_tolerance=6,
    snap_tolerance=4,
    join_tolerance=4,
    edge_min_length=18,
    min_words_vertical=1,
    min_words_horizontal=1,
)

_PDF_STREAM = dict(
    vertical_strategy="text",
    horizontal_strategy="text",
    text_x_tolerance=2,
    text_y_tolerance=2,
)


# ── Header promotion ─────────────────────────────────────────────────────────

def _promote_header(df: pd.DataFrame) -> pd.DataFrame:
    """If column headers look auto-generated, promote the first row."""
    if df is None or df.empty:
        return df

    def looks_default(cols):
        try:
            as_int = [int(str(c)) for c in cols]
            if as_int == list(range(len(cols))):
                return True
        except Exception:
            pass
        return all(str(c).startswith("Unnamed") for c in cols)

    if not looks_default(df.columns):
        return df

    first = df.iloc[0].astype(str).str.strip().tolist()
    non_empty = sum(bool(x) for x in first)
    if non_empty >= max(2, int(0.6 * len(first))):
        seen, new_cols = {}, []
        for x in first:
            base = x if x else "Column"
            cnt  = seen.get(base, 0)
            seen[base] = cnt + 1
            new_cols.append(base if cnt == 0 else f"{base}_{cnt + 1}")
        df = df.iloc[1:].copy()
        df.columns = new_cols
        df.reset_index(drop=True, inplace=True)
    return df


# ── Fake-table filter ─────────────────────────────────────────────────────────

def _is_fake_table(df: pd.DataFrame) -> bool:
    """Return True if the DataFrame is not a real data table."""
    if df is None or df.empty:
        return True
    rows, cols = df.shape
    if rows < 2 or cols < 2:
        return True
    empty_ratio = df.isna().sum().sum() / (rows * cols)
    if empty_ratio > 0.70:
        return True
    flat = " ".join(df.astype(str).fillna("").values.flatten()).lower()
    suspicious = ["see ", "page ", "none", "n/a", "refer", "note ", "footnote", "respecting nature"]
    if any(m in flat for m in suspicious):
        return True
    non_empty_cols = (df.astype(str).apply(lambda s: s.str.strip().ne("")).sum(axis=0) > 1).sum()
    if non_empty_cols < 2:
        return True
    return False


# ── Clean DataFrame ───────────────────────────────────────────────────────────

def _clean_df(df: pd.DataFrame) -> pd.DataFrame:
    """Drop fully empty rows/columns."""
    df = df.copy()
    df.replace({None: ""}, inplace=True)
    df = df.loc[:, (df.astype(str).ne("")).any(axis=0)]
    df = df.loc[(df.astype(str).ne("")).any(axis=1)]
    return df


# ── Per-format extractors ─────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes, filename: str) -> List[Dict]:
    if not PDF_AVAILABLE:
        return []
    tables = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for pno, page in enumerate(pdf.pages, start=1):
            found  = list(page.find_tables(table_settings=_PDF_LATTICE) or [])
            found += list(page.find_tables(table_settings=_PDF_STREAM)   or [])
            if not found:
                continue

            seen_sigs, tidx = set(), 0
            for t in found:
                try:
                    data = t.extract()
                    df   = pd.DataFrame(data)
                except Exception:
                    continue

                df  = _clean_df(df)
                sig = "|".join(
                    "\t".join(map(str, df.iloc[i, :].tolist()[:8]))
                    for i in range(min(len(df), 3))
                )
                if sig in seen_sigs:
                    continue
                seen_sigs.add(sig)

                df = _promote_header(df)
                if _is_fake_table(df):
                    continue

                tables.append({
                    "filename":  filename,
                    "table_name": f"{filename}_p{pno}_t{tidx}",
                    "source":    f"Page {pno}, Table {tidx + 1}",
                    "df":        df,
                    "row_count": len(df),
                    "col_count": len(df.columns),
                })
                tidx += 1
    return tables


def _extract_docx(file_bytes: bytes, filename: str) -> List[Dict]:
    if not DOCX_AVAILABLE:
        return []
    tables = []
    doc = DocxDocument(io.BytesIO(file_bytes))
    for t_i, table in enumerate(doc.tables):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        df = _clean_df(pd.DataFrame(rows))
        df = _promote_header(df)
        if _is_fake_table(df):
            continue
        tables.append({
            "filename":   filename,
            "table_name": f"{filename}_t{t_i}",
            "source":     f"Table {t_i + 1}",
            "df":         df,
            "row_count":  len(df),
            "col_count":  len(df.columns),
        })
    return tables


def _extract_pptx(file_bytes: bytes, filename: str) -> List[Dict]:
    if not PPTX_AVAILABLE:
        return []
    tables, count = [], 0
    pres = Presentation(io.BytesIO(file_bytes))
    for s_i, slide in enumerate(pres.slides):
        for shape in slide.shapes:
            if not hasattr(shape, "table"):
                continue
            rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
            if not rows:
                continue
            df = _clean_df(pd.DataFrame(rows))
            df = _promote_header(df)
            if _is_fake_table(df):
                continue
            tables.append({
                "filename":   filename,
                "table_name": f"{filename}_s{s_i}_t{count}",
                "source":     f"Slide {s_i + 1}, Table {count + 1}",
                "df":         df,
                "row_count":  len(df),
                "col_count":  len(df.columns),
            })
            count += 1
    return tables


def _extract_excel(file_bytes: bytes, filename: str) -> List[Dict]:
    tables = []
    try:
        xl = pd.ExcelFile(io.BytesIO(file_bytes))
        for sheet in xl.sheet_names:
            df = xl.parse(sheet).reset_index(drop=True)
            df = _clean_df(df)
            df = _promote_header(df)
            if _is_fake_table(df):
                continue
            tables.append({
                "filename":   filename,
                "table_name": f"{filename}_{sheet}",
                "source":     f"Sheet: {sheet}",
                "df":         df,
                "row_count":  len(df),
                "col_count":  len(df.columns),
            })
    except Exception as e:
        print(f"[table_extractor] Excel error {filename}: {e}")
    return tables


def _extract_csv(file_bytes: bytes, filename: str) -> List[Dict]:
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), encoding="utf-8", engine="python")
        df = _clean_df(df)
        df = _promote_header(df)
        if _is_fake_table(df):
            return []
        return [{
            "filename":   filename,
            "table_name": filename,
            "source":     "CSV file",
            "df":         df,
            "row_count":  len(df),
            "col_count":  len(df.columns),
        }]
    except Exception as e:
        print(f"[table_extractor] CSV error {filename}: {e}")
        return []


# ── Public: extract tables from an uploaded file ──────────────────────────────

def extract_tables_from_upload(uploaded_file) -> List[Dict]:
    """
    Extract all tables from a Streamlit UploadedFile object.
    Returns a list of table dicts, each with: filename, table_name,
    source, df (DataFrame), row_count, col_count.
    """
    filename  = uploaded_file.name
    file_ext  = filename.rsplit(".", 1)[-1].lower()
    file_bytes = uploaded_file.getvalue()   # bytes, doesn't consume the stream

    dispatch = {
        "pdf":  _extract_pdf,
        "docx": _extract_docx,
        "doc":  _extract_docx,
        "pptx": _extract_pptx,
        "xlsx": _extract_excel,
        "xls":  _extract_excel,
        "csv":  _extract_csv,
    }

    fn = dispatch.get(file_ext)
    if fn is None:
        return []

    try:
        return fn(file_bytes, filename)
    except Exception as e:
        print(f"[table_extractor] Failed on {filename}: {e}")
        return []


def batch_extract_tables(uploaded_files) -> List[Dict]:
    """Extract tables from all uploaded files, return flat list."""
    all_tables = []
    for f in uploaded_files:
        all_tables.extend(extract_tables_from_upload(f))
    return all_tables


# ── Convert tables → text chunks (for search / analysis / RAG) ───────────────

def table_to_text(table: Dict, max_rows: int = 50) -> str:
    """
    Render a table dict as a plain-text representation for the LLM.
    Format is intentionally compact so it fits in context windows.
    """
    df      = table["df"]
    source  = table.get("source", "")
    fname   = table.get("filename", "")
    cols    = list(df.columns)
    n_rows  = len(df)

    lines = [
        f"TABLE from {fname} ({source})",
        f"Columns ({len(cols)}): {', '.join(str(c) for c in cols)}",
        f"Rows: {n_rows}",
        "---",
    ]
    display_df = df.head(max_rows)
    for _, row in display_df.iterrows():
        row_str = " | ".join(f"{c}: {v}" for c, v in zip(cols, row.tolist()))
        lines.append(row_str)
    if n_rows > max_rows:
        lines.append(f"... ({n_rows - max_rows} more rows not shown)")
    return "\n".join(lines)


def tables_to_chunks(tables: List[Dict]) -> List[Dict]:
    """
    Convert extracted table list into chunk dicts compatible with the
    existing chunker/search/analyzer pipeline.
    Each table becomes one chunk tagged with is_table=True.
    """
    chunks = []
    for i, table in enumerate(tables):
        text = table_to_text(table)
        chunks.append({
            "chunk_id":    f"{table['filename']}::table_{i}",
            "filename":    table["filename"],
            "chunk_index": i,
            "start_line":  1,
            "end_line":    1,
            "text":        text,
            "char_count":  len(text),
            "word_count":  len(text.split()),
            "is_table":    True,
            "table_name":  table["table_name"],
            "source":      table.get("source", ""),
        })
    return chunks